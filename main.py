#!/usr/bin/python3

import eel
import unicodedata
import json
import base64
import shutil
import time
from pathlib import Path
from PIL import Image
import io
from docx import Document
from docx.shared import Inches, Cm
from pptx import Presentation
from pptx.util import Inches as PptxInches
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import mm
from reportlab.lib.enums import TA_CENTER
from screeninfo import get_monitors

# ---------- Transliteration mapping ----------
base_consonants = {
    'ሀ': 'h', 'ለ': 'l', 'ሐ': 'h', 'መ': 'm', 'ሠ': 's',
    'ረ': 'r', 'ሰ': 's', 'ሸ': 'sh', 'ቀ': 'q', 'በ': 'b',
    'ተ': 't', 'ቸ': 'ch', 'ኀ': 'h', 'ነ': 'n', 'ኘ': 'ny',
    'አ': 'ə',
    'ከ': 'k', 'ኸ': 'h', 'ወ': 'w', 'ዐ': 'ə',
    'ዘ': 'z', 'ዠ': 'zh', 'የ': 'y', 'ደ': 'd',
    'ጀ': 'j', 'ገ': 'g', 'ጠ': 't', 'ጨ': 'ch',
    'ጰ': 'p', 'ጸ': 'ts', 'ፀ': 'ts', 'ፈ': 'f', 'ፐ': 'p',
}

vowel_suffixes = ['a', 'u', 'i', 'a', 'e', '', 'o']

amharic_to_sound = {}
for first_order_char, cons in base_consonants.items():
    base_code = ord(first_order_char)
    for order in range(7):
        char_code = base_code + order
        if 0x1200 <= char_code <= 0x137F:
            syllable = chr(char_code)
            try:
                name = unicodedata.name(syllable)
                if 'ETHIOPIC' not in name:
                    continue
            except (ValueError, TypeError):
                continue
            sound = cons + vowel_suffixes[order]
            amharic_to_sound[syllable] = sound

extra_mapping = {
    'ቈ': 'qwa', 'ቊ': 'qwi', 'ቋ': 'qwa', 'ቌ': 'qwe', 'ቍ': 'qw',
    'ኈ': 'hwa', 'ኊ': 'hwi', 'ኋ': 'hwa', 'ኌ': 'hwe', 'ኍ': 'hw',
    'ዀ': 'hwa', 'ዂ': 'hwi', 'ዃ': 'hwa', 'ዄ': 'hwe',
    'ጐ': 'gwa', 'ጒ': 'gwi', 'ጓ': 'gwa', 'ጔ': 'gwe', 'ጕ': 'gw',
}
amharic_to_sound.update(extra_mapping)

@eel.expose
def amharic_to_english_sound(text: str) -> str:
    result = []
    for ch in text:
        result.append(amharic_to_sound.get(ch, ch))
    return ''.join(result)

# ---------- Project Management ----------
PROJECTS_DIR = Path(__file__).parent / "projects"
PROJECTS_DIR.mkdir(exist_ok=True)

def _get_project_path(name: str) -> Path:
    if not name:
        raise ValueError("Project name cannot be empty or None")
    safe_name = "".join(c for c in str(name) if c.isalnum() or c in " ._-")
    if not safe_name.strip():
        safe_name = "unnamed"
    project_path = PROJECTS_DIR / safe_name
    project_path.mkdir(exist_ok=True)
    return project_path

def _get_metadata_path(project_path: Path) -> Path:
    return project_path / "metadata.json"

def _load_metadata(project_path: Path) -> dict:
    meta_path = _get_metadata_path(project_path)
    if meta_path.exists():
        with open(meta_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        if "last_modified" not in data:
            data["last_modified"] = 0
        return data
    return {"converted_text": "", "photos": [], "last_modified": 0}

def _save_metadata(project_path: Path, data: dict):
    meta_path = _get_metadata_path(project_path)
    data["last_modified"] = time.time()
    with open(meta_path, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

@eel.expose
def list_projects():
    projects = []
    for item in PROJECTS_DIR.iterdir():
        if item.is_dir() and _get_metadata_path(item).exists():
            projects.append(item.name)
    return projects

@eel.expose
def get_all_projects():
    """Return list of all project names with metadata for merge tab."""
    projects = []
    for item in PROJECTS_DIR.iterdir():
        if item.is_dir() and _get_metadata_path(item).exists():
            metadata = _load_metadata(item)
            projects.append({
                "name": item.name,
                "text_length": len(metadata.get("converted_text", "")),
                "photo_count": len(metadata.get("photos", []))
            })
    return projects

@eel.expose
def create_project(name: str) -> bool:
    if not name:
        return False
    try:
        project_path = _get_project_path(name)
        if _get_metadata_path(project_path).exists():
            return False
        _save_metadata(project_path, {"converted_text": "", "photos": []})
        return True
    except Exception:
        return False

@eel.expose
def delete_project(project_name: str) -> bool:
    if not project_name:
        return False
    try:
        project_path = _get_project_path(project_name)
        if project_path.exists():
            shutil.rmtree(project_path)
            return True
        return False
    except Exception:
        return False

@eel.expose
def get_project_data(project_name: str):
    if not project_name:
        return {"converted_text": "", "photos": []}
    try:
        project_path = _get_project_path(project_name)
        metadata = _load_metadata(project_path)
        converted_text = metadata.get("converted_text", "")
        photos = []
        for photo_rel in metadata.get("photos", []):
            photo_path = project_path / photo_rel
            if photo_path.exists():
                try:
                    with Image.open(photo_path) as img:
                        img.thumbnail((200, 200))
                        buffered = io.BytesIO()
                        img.save(buffered, format="JPEG" if img.format != 'PNG' else "PNG")
                        img_base64 = base64.b64encode(buffered.getvalue()).decode('utf-8')
                        mime = "image/jpeg" if img.format != 'PNG' else "image/png"
                        photos.append({
                            "name": photo_rel,
                            "thumb": f"data:{mime};base64,{img_base64}"
                        })
                except Exception:
                    photos.append({"name": photo_rel, "thumb": ""})
            else:
                photos.append({"name": photo_rel, "thumb": ""})
        return {"converted_text": converted_text, "photos": photos}
    except Exception:
        return {"converted_text": "", "photos": []}

@eel.expose
def upload_photos(project_name: str, photo_data_list: list):
    if not project_name:
        return []
    try:
        project_path = _get_project_path(project_name)
        metadata = _load_metadata(project_path)
        saved_names = []
        for item in photo_data_list:
            filename = item.get('name', '')
            if not filename:
                continue
            safe_filename = "".join(c for c in filename if c.isalnum() or c in "._- ")
            if not safe_filename:
                continue
            data_url = item.get('data', '')
            if ',' in data_url:
                base64_data = data_url.split(',', 1)[1]
            else:
                base64_data = data_url
            try:
                img_data = base64.b64decode(base64_data)
                file_path = project_path / safe_filename
                counter = 1
                original_path = file_path
                while file_path.exists():
                    stem = original_path.stem
                    suffix = original_path.suffix
                    file_path = original_path.parent / f"{stem}_{counter}{suffix}"
                    counter += 1
                with open(file_path, 'wb') as f:
                    f.write(img_data)
                saved_names.append(file_path.name)
            except Exception:
                continue

        existing_photos = metadata.get("photos", [])
        for name in saved_names:
            if name not in existing_photos:
                existing_photos.append(name)
        metadata["photos"] = existing_photos
        _save_metadata(project_path, metadata)
        return saved_names
    except Exception:
        return []

@eel.expose
def save_converted_text(project_name: str, text: str):
    if not project_name:
        return False
    try:
        project_path = _get_project_path(project_name)
        metadata = _load_metadata(project_path)
        metadata["converted_text"] = text
        _save_metadata(project_path, metadata)
        return True
    except Exception:
        return False

@eel.expose
def export_project_to_docx(project_name: str) -> str:
    if not project_name:
        return ""
    try:
        project_path = _get_project_path(project_name)
        metadata = _load_metadata(project_path)
        doc = Document()
        doc.add_heading(f"Project: {project_name}", 0)
        doc.add_heading("Converted Text (Latin Harari)", level=1)
        doc.add_paragraph(metadata.get("converted_text", "") or "(No text saved)")
        photos = metadata.get("photos", [])
        if photos:
            doc.add_heading("Photos", level=1)
            cols = 3
            rows = (len(photos) + cols - 1) // cols
            table = doc.add_table(rows=rows, cols=cols)
            table.style = 'Table Grid'
            table.autofit = False
            for cell in table.columns:
                cell.width = Cm(5.5)
            for idx, photo_rel in enumerate(photos):
                row = idx // cols
                col = idx % cols
                cell = table.cell(row, col)
                cell.paragraphs[0].clear()
                photo_path = project_path / photo_rel
                if photo_path.exists():
                    try:
                        run = cell.paragraphs[0].add_run()
                        run.add_picture(str(photo_path), width=Cm(5.0))
                        cell.add_paragraph(photo_rel, style='Caption')
                    except Exception:
                        cell.text = f"[Could not embed: {photo_rel}]"
                else:
                    cell.text = f"[Missing: {photo_rel}]"
        else:
            doc.add_paragraph("No photos attached to this project.")
        export_path = project_path / f"{project_name}_export.docx"
        doc.save(str(export_path))
        return str(export_path)
    except Exception:
        return ""

@eel.expose
def export_project_to_pptx(project_name: str) -> str:
    """Export project to PowerPoint presentation."""
    if not project_name:
        return ""
    try:
        project_path = _get_project_path(project_name)
        metadata = _load_metadata(project_path)
        prs = Presentation()
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = f"Project: {project_name}"
        subtitle.text = "Exported from Amharic Converter"
        # Text slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        title.text = "Converted Text (Latin Harari)"
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = metadata.get("converted_text", "") or "(No text saved)"
        # Photos slides
        photos = metadata.get("photos", [])
        if photos:
            photos_per_slide = 6
            for i in range(0, len(photos), photos_per_slide):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                left = PptxInches(0.5)
                top = PptxInches(1.5)
                width = PptxInches(3)
                height = PptxInches(2.5)
                for j, photo_rel in enumerate(photos[i:i+photos_per_slide]):
                    photo_path = project_path / photo_rel
                    if photo_path.exists():
                        try:
                            col = j % 3
                            row = j // 3
                            pic_left = left + (col * (width + PptxInches(0.3)))
                            pic_top = top + (row * (height + PptxInches(0.5)))
                            slide.shapes.add_picture(str(photo_path), pic_left, pic_top, width=width, height=height)
                        except Exception:
                            pass
        export_path = project_path / f"{project_name}_export.pptx"
        prs.save(str(export_path))
        return str(export_path)
    except Exception as e:
        print(f"PPTX export error: {e}")
        return ""

@eel.expose
def export_project_to_pdf(project_name: str) -> str:
    """Export project to PDF using ReportLab."""
    if not project_name:
        return ""
    try:
        project_path = _get_project_path(project_name)
        metadata = _load_metadata(project_path)
        pdf_path = project_path / f"{project_name}_export.pdf"
        doc = SimpleDocTemplate(str(pdf_path), pagesize=A4,
                                rightMargin=20*mm, leftMargin=20*mm,
                                topMargin=20*mm, bottomMargin=20*mm)
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle('CustomTitle', parent=styles['Title'],
                                     alignment=TA_CENTER, fontSize=16, spaceAfter=12)
        heading_style = ParagraphStyle('CustomHeading', parent=styles['Heading1'],
                                       fontSize=14, spaceAfter=8)
        normal_style = styles['Normal']
        story = []
        story.append(Paragraph(f"Project: {project_name}", title_style))
        story.append(Spacer(1, 12))
        story.append(Paragraph("Converted Text (Latin Harari)", heading_style))
        story.append(Spacer(1, 6))
        text = metadata.get("converted_text", "") or "(No text saved)"
        story.append(Paragraph(text.replace('\n', '<br/>'), normal_style))
        story.append(Spacer(1, 12))
        photos = metadata.get("photos", [])
        if photos:
            story.append(Paragraph("Photos", heading_style))
            story.append(Spacer(1, 6))
            for idx, photo_rel in enumerate(photos):
                photo_path = project_path / photo_rel
                if photo_path.exists():
                    try:
                        img = RLImage(str(photo_path), width=80*mm, height=80*mm)
                        story.append(img)
                        story.append(Paragraph(photo_rel, normal_style))
                        if (idx + 1) % 3 == 0:
                            story.append(PageBreak())
                        else:
                            story.append(Spacer(1, 6))
                    except Exception:
                        story.append(Paragraph(f"[Could not embed: {photo_rel}]", normal_style))
        doc.build(story)
        return str(pdf_path)
    except Exception as e:
        print(f"PDF export error: {e}")
        return ""

@eel.expose
def merge_projects(project_names: list, new_project_name: str, export_format: str) -> dict:
    if not project_names or not new_project_name:
        return {"success": False, "error": "No projects selected or invalid name"}
    try:
        success = create_project(new_project_name)
        if not success:
            return {"success": False, "error": "Project name already exists or invalid"}
        new_project_path = _get_project_path(new_project_name)
        merged_text = ""
        all_photos = []
        for proj_name in project_names:
            proj_path = _get_project_path(proj_name)
            metadata = _load_metadata(proj_path)
            merged_text += f"\n\n--- From project: {proj_name} ---\n\n"
            merged_text += metadata.get("converted_text", "") + "\n"
            for photo_rel in metadata.get("photos", []):
                src_photo = proj_path / photo_rel
                if src_photo.exists():
                    new_name = f"{proj_name}_{photo_rel}"
                    dst_photo = new_project_path / new_name
                    counter = 1
                    original_dst = dst_photo
                    while dst_photo.exists():
                        stem = original_dst.stem
                        suffix = original_dst.suffix
                        dst_photo = original_dst.parent / f"{stem}_{counter}{suffix}"
                        counter += 1
                    shutil.copy2(src_photo, dst_photo)
                    all_photos.append(dst_photo.name)
        merged_text = merged_text.strip()
        _save_metadata(new_project_path, {"converted_text": merged_text, "photos": all_photos})
        export_path = ""
        if export_format == "docx":
            export_path = export_project_to_docx(new_project_name)
        elif export_format == "pptx":
            export_path = export_project_to_pptx(new_project_name)
        elif export_format == "pdf":
            export_path = export_project_to_pdf(new_project_name)
        if not export_path:
            return {"success": False, "error": f"Export to {export_format} failed"}
        return {"success": True, "new_project_name": new_project_name, "export_path": export_path}
    except Exception as e:
        return {"success": False, "error": str(e)}

@eel.expose
def get_dashboard_stats():
    total_projects = 0
    total_photos = 0
    total_chars = 0
    recent = []
    for item in PROJECTS_DIR.iterdir():
        if item.is_dir() and _get_metadata_path(item).exists():
            metadata = _load_metadata(item)
            total_projects += 1
            total_photos += len(metadata.get("photos", []))
            total_chars += len(metadata.get("converted_text", ""))
            recent.append({
                "name": item.name,
                "last_modified": metadata.get("last_modified", 0),
                "photo_count": len(metadata.get("photos", [])),
                "text_length": len(metadata.get("converted_text", ""))
            })
    recent.sort(key=lambda x: x["last_modified"], reverse=True)
    return {
        "total_projects": total_projects,
        "total_photos": total_photos,
        "total_characters": total_chars,
        "recent_projects": recent[:5]
    }

# ---------- Start the app ----------
eel.init('ui')
try:
    screen = get_monitors()[0]
    screen_width = screen.width - screen.width * 0.1
    screen_height = screen.height - screen.height * 0.1
except:
    screen_width = 1200
    screen_height = 800
eel.start('index.html', size=(int(screen_width), int(screen_height)))
